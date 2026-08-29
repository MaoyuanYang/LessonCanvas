"""PPTX lesson-deck rendering and structural validation tools.

Registered with MCP-compatible definitions (ADR-0004) and consumed by the
deck workflow. Generated model content is untrusted input: it is rendered as
inert slide/notes text and can never alter tool selection or policy.
"""

import io

from pptx import Presentation

OBJECTIVES_SLIDE_TITLE = "教学目标"
KEY_POINTS_SLIDE_TITLE = "重点与难点"
STAGE_SLIDE_TITLE_PREFIX = "教学过程"
HOMEWORK_SLIDE_TITLE = "作业布置"

REQUIRED_SLIDE_TITLES = (OBJECTIVES_SLIDE_TITLE, KEY_POINTS_SLIDE_TITLE, HOMEWORK_SLIDE_TITLE)

MIN_SLIDE_COUNT = 5  # title + objectives + key points + >=1 stage + homework
MIN_DECK_TEXT_CHARS = 30

RENDER_TOOL_DEFINITION = {
    "name": "render_lesson_deck_pptx",
    "description": (
        "Render one structured lesson deck into an editable PPTX file with the "
        "standard senior-high slide grammar: title, objectives, key points, "
        "bounded teaching-stage slides, and homework; teacher notes and "
        "citations go to speaker notes."
    ),
    "inputSchema": {
        "type": "object",
        "properties": {
            "slide_deck": {
                "type": "object",
                "description": "structured deck content from the deck writer specialist",
            },
            "lesson_index": {"type": "integer"},
            "language_mode": {"type": "string"},
        },
        "required": ["slide_deck", "lesson_index", "language_mode"],
    },
}

VALIDATE_TOOL_DEFINITION = {
    "name": "validate_lesson_deck_pptx",
    "description": (
        "Structurally validate a rendered lesson-deck PPTX: openable, required "
        "slides present, every slide carries non-empty editable text frames "
        "(not whole-slide images), slide count within configured bounds. Does "
        "not judge content quality."
    ),
    "inputSchema": {
        "type": "object",
        "properties": {"content_b64": {"type": "string", "description": "rendered file bytes"}},
        "required": ["content_b64"],
    },
}


def _clean_lines(values) -> list[str]:
    return [str(item).strip() for item in (values or []) if str(item).strip()]


def _add_content_slide(
    presentation: Presentation, title: str, bullets: list[str], notes: str | None
) -> None:
    slide = presentation.slides.add_slide(presentation.slide_layouts[1])
    slide.shapes.title.text = title
    text_frame = slide.placeholders[1].text_frame
    text_frame.word_wrap = True
    lines = bullets or ["（待补充）"]
    text_frame.text = lines[0]
    for line in lines[1:]:
        paragraph = text_frame.add_paragraph()
        paragraph.text = line
        paragraph.level = 0
    if notes:
        slide.notes_slide.notes_text_frame.text = notes


def render_lesson_deck_pptx(slide_deck: dict, lesson_index: int, language_mode: str) -> bytes:
    deck = slide_deck if isinstance(slide_deck, dict) else {}
    presentation = Presentation()

    title = str(deck.get("title") or f"第{lesson_index}课")[:120]
    subtitle = str(deck.get("unit_title") or "").strip()[:120]
    notes_lines = _clean_lines(deck.get("notes"))
    notes_text = "\n".join(notes_lines) if notes_lines else None

    title_slide = presentation.slides.add_slide(presentation.slide_layouts[0])
    title_slide.shapes.title.text = f"第{lesson_index}课 {title}"
    if subtitle and 1 in title_slide.placeholders:
        title_slide.placeholders[1].text_frame.text = subtitle
    if notes_text:
        title_slide.notes_slide.notes_text_frame.text = notes_text

    _add_content_slide(
        presentation,
        OBJECTIVES_SLIDE_TITLE,
        _clean_lines(deck.get("objectives")),
        None,
    )

    key_lines = [f"重点：{item}" for item in _clean_lines(deck.get("key_points"))]
    key_lines += [f"难点：{item}" for item in _clean_lines(deck.get("difficulties"))]
    _add_content_slide(presentation, KEY_POINTS_SLIDE_TITLE, key_lines, None)

    stage_number = 0
    for stage in deck.get("stage_slides") or []:
        if not isinstance(stage, dict):
            continue
        stage_number += 1
        heading = str(stage.get("heading") or f"{STAGE_SLIDE_TITLE_PREFIX}（{stage_number}）")
        _add_content_slide(
            presentation,
            heading,
            _clean_lines(stage.get("bullets")),
            None,
        )

    _add_content_slide(
        presentation,
        HOMEWORK_SLIDE_TITLE,
        [str(deck.get("homework") or "（待补充）")],
        None,
    )

    buffer = io.BytesIO()
    presentation.save(buffer)
    return buffer.getvalue()


def slide_count_of(content: bytes) -> int:
    try:
        return len(Presentation(io.BytesIO(content)).slides)
    except Exception:  # noqa: BLE001 - unreachable after successful validation
        return 0


def _slide_texts(slide) -> list[str]:
    texts: list[str] = []
    for shape in slide.shapes:
        if getattr(shape, "has_text_frame", False):
            text = shape.text_frame.text.strip()
            if text:
                texts.append(text)
    return texts


def validate_lesson_deck_pptx(
    content: bytes, max_slides: int | None = None
) -> tuple[bool, str | None]:
    """Return (ok, reason). Structural checks only (Spec D7).

    Editability evidence: python-pptx only reads text from text frames, so a
    slide whose only content is a picture contributes no text and fails the
    per-slide non-empty check.
    """

    from lessoncanvas.settings import get_settings

    if not content:
        return False, "empty file"
    try:
        presentation = Presentation(io.BytesIO(content))
    except Exception as error:  # noqa: BLE001 - any parse failure is a validation failure
        return False, f"unopenable: {type(error).__name__}"

    limit = max_slides if max_slides is not None else get_settings().deck_max_slides
    slides = list(presentation.slides)
    if len(slides) < MIN_SLIDE_COUNT:
        return False, f"too few slides: {len(slides)} < {MIN_SLIDE_COUNT}"
    if len(slides) > limit:
        return False, f"too many slides: {len(slides)} > {limit}"

    all_titles: set[str] = set()
    total_chars = 0
    for index, slide in enumerate(slides):
        texts = _slide_texts(slide)
        if not texts:
            return False, f"slide {index + 1} has no editable text"
        all_titles.update(texts)
        total_chars += sum(len(text) for text in texts)

    missing = [
        title
        for title in REQUIRED_SLIDE_TITLES
        if not any(title in text for text in all_titles)
    ]
    if missing:
        return False, f"missing slides: {','.join(missing)}"
    has_stage_slide = any(
        text.startswith(STAGE_SLIDE_TITLE_PREFIX)
        for text in all_titles
        if text not in REQUIRED_SLIDE_TITLES
    )
    if not has_stage_slide:
        return False, "missing stage slide"
    if total_chars < MIN_DECK_TEXT_CHARS:
        return False, "deck text is empty"
    return True, None


def execute_pptx_tool(name: str, arguments: dict) -> dict:
    import base64

    if name == RENDER_TOOL_DEFINITION["name"]:
        content = render_lesson_deck_pptx(
            arguments["slide_deck"],
            int(arguments["lesson_index"]),
            str(arguments["language_mode"]),
        )
        return {"content_b64": base64.b64encode(content).decode(), "size_bytes": len(content)}
    if name == VALIDATE_TOOL_DEFINITION["name"]:
        content = base64.b64decode(arguments["content_b64"])
        ok, reason = validate_lesson_deck_pptx(content)
        return {"ok": ok, "reason": reason}
    raise KeyError(f"unknown tool: {name}")
